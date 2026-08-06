#!/usr/bin/env python3
"""KamuKOD/KamuYZ sunum üretici — MD → PPTX + HTML deck"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os, sys

OUT = "/Users/murat/Desktop/youtube/kamuyz-agent-starter/pazarlama/sunumlar/pptx"
os.makedirs(OUT, exist_ok=True)

BG = RGBColor(0x0D, 0x0D, 0x0D)
ACCENT = RGBColor(0xD4, 0x85, 0x6A)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)

W = Inches(13.333)
H = Inches(7.5)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_accent_bar(slide, left=Inches(2), top=Inches(0.3), width=Inches(0.8), height=Inches(0.04)):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def cover(prs, kicker, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    if kicker:
        tb = slide.shapes.add_textbox(Inches(2), Inches(1.8), Inches(9), Inches(0.5))
        tf = tb.text_frame
        tf.paragraphs[0].text = kicker
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = ACCENT
        tf.paragraphs[0].font.name = 'Space Mono'

    add_accent_bar(slide, Inches(2), Inches(2.5))

    tb = slide.shapes.add_textbox(Inches(2), Inches(2.8), Inches(10), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for line in title.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(56)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Outfit'
        p.space_after = Pt(4)

    if subtitle:
        tb = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(10), Inches(0.8))
        tf = tb.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.italic = True
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.paragraphs[0].font.name = 'Outfit'

    return slide

def quote_slide(prs, quote, author=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_accent_bar(slide, Inches(5.8), Inches(2.8), Inches(1.6), Inches(0.04))

    tb = slide.shapes.add_textbox(Inches(2), Inches(3.2), Inches(9.3), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = f'"{quote}"'
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.italic = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = 'Outfit'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if author:
        tb2 = slide.shapes.add_textbox(Inches(2), Inches(6), Inches(9.3), Inches(0.5))
        tf2 = tb2.text_frame
        tf2.paragraphs[0].text = author
        tf2.paragraphs[0].font.size = Pt(14)
        tf2.paragraphs[0].font.color.rgb = ACCENT
        tf2.paragraphs[0].font.name = 'Space Mono'
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def title_body(prs, kicker, title, body_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    if kicker:
        tb = slide.shapes.add_textbox(Inches(2), Inches(1.2), Inches(9), Inches(0.5))
        tf = tb.text_frame
        tf.paragraphs[0].text = kicker
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = ACCENT
        tf.paragraphs[0].font.name = 'Space Mono'

    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(2), Inches(1.9), Inches(10), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for line in title.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Outfit'
        p.space_after = Pt(4)

    tb2 = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(10), Inches(3.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for line in body_lines:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.font.name = 'Outfit'
        p.space_after = Pt(8)

    return slide

def compare_slide(prs, kicker, title, left_title, left_items, right_title, right_items, left_is_bad=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    if kicker:
        tb = slide.shapes.add_textbox(Inches(2), Inches(1.2), Inches(9), Inches(0.5))
        tf = tb.text_frame
        tf.paragraphs[0].text = kicker
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.color.rgb = ACCENT
        tf.paragraphs[0].font.name = 'Space Mono'

    add_accent_bar(slide)

    tb = slide.shapes.add_textbox(Inches(2), Inches(1.9), Inches(10), Inches(1))
    tf = tb.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.name = 'Outfit'

    # Left
    lshape = slide.shapes.add_shape(1, Inches(2), Inches(3.2), Inches(4.5), Inches(3.8))
    lshape.fill.solid(); lshape.fill.fore_color.rgb = RGBColor(0x1a,0x1a,0x1a)
    lshape.line.color.rgb = RED if left_is_bad else GREEN

    tb = slide.shapes.add_textbox(Inches(2.2), Inches(3.4), Inches(4.1), Inches(0.5))
    tf = tb.text_frame
    tf.paragraphs[0].text = left_title
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RED if left_is_bad else GREEN
    tf.paragraphs[0].font.name = 'Outfit'

    tb2 = slide.shapes.add_textbox(Inches(2.2), Inches(4), Inches(4.1), Inches(3))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for line in left_items:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.font.name = 'Outfit'
        p.space_after = Pt(6)

    # Right
    rshape = slide.shapes.add_shape(1, Inches(7), Inches(3.2), Inches(4.5), Inches(3.8))
    rshape.fill.solid(); rshape.fill.fore_color.rgb = RGBColor(0x1a,0x1a,0x1a)
    rshape.line.color.rgb = GREEN if left_is_bad else RED

    tb = slide.shapes.add_textbox(Inches(7.2), Inches(3.4), Inches(4.1), Inches(0.5))
    tf = tb.text_frame
    tf.paragraphs[0].text = right_title
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = GREEN if left_is_bad else RED
    tf.paragraphs[0].font.name = 'Outfit'

    tb2 = slide.shapes.add_textbox(Inches(7.2), Inches(4), Inches(4.1), Inches(3))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for line in right_items:
        p = tf2.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY
        p.font.name = 'Outfit'
        p.space_after = Pt(6)

    return slide

# === HERMES INTRO DECK ===
print("[1/3] Hermes Intro Deck...")
prs = Presentation()
prs.slide_width = W; prs.slide_height = H

cover(prs, "Hermes Agent", "Kendi Kendine Öğrenen\nYapay Zeka Ajanı", "Nous Research · MIT Lisans · Açık Kaynak")
title_body(prs, "Diğerlerinden Farkı", "Neden Hermes?",
    ["ChatGPT: Soru sor, cevap al",
     "OpenClaw: Görev ver, yapar (her seferinde sıfırdan)",
     "Hermes: Görev ver, yapar, nasıl yaptığını kaydeder, bir dahaki sefere daha iyi yapar",
     "",
     "Zamanla sana özel bir uzman haline gelir. Ne kadar kullanırsan, o kadar iyi olur."])
quote_slide(prs, "Hermes'e geçince aynı iş akışı ilk denemede sorunsuz çalıştı, bir hafta boyunca müdahalesiz devam etti.", "Reddit · 342 upvote")
title_body(prs, "Mimari", "4 Katmanlı Hafıza",
    ["L1 — Session Context: Anlık konuşma, oturum kapanınca silinir",
     "L2 — Persisted Facts: ~2.200 karakter (MEMORY.md), sen silene kadar kalır",
     "L3 — Session Archive: SQLite FTS5, tüm geçmiş oturumlar tam metin aramayla",
     "L4 — Procedural Skills: Kendi yazdığı SKILL.md, ihtiyaç anında dinamik yüklenir"])
title_body(prs, "Öğrenme Döngüsü", "Kapalı Devre (Closed Loop)",
    ["Görev → Tamamla → Analiz → SKILL.md oluştur → Kullandıkça iyileştir",
     "",
     "Bu döngü sayesinde ajan her görevden sonra biraz daha iyi olur.",
     "Başarılı iş akışlarını kendi SKILL.md dosyasına yazar, tekrar kullanır."])
compare_slide(prs, "Karşılaştırma", "OpenClaw vs Hermes",
    "OpenClaw", ["50+ kanal", "13.000+ hazır skill", "TypeScript/Node.js", "ClawHub marketplace", "CVE-2026-25253 yaşadı"],
    "Hermes Agent", ["Kendi skill'ini yazar", "4 katman hafıza", "Python (%88)", "Öğrenme döngüsü", "Sıfır bildirilen CVE"],
    left_is_bad=False)
title_body(prs, "Kurulum", "25 dakikada çalışır halde",
    ["OpenClaw: ~4 saat (kullanıcı deneyimi)",
     "Hermes: ~25 dakika (kullanıcı deneyimi)",
     "",
     "Bizim güvenli paketimiz: Tek komutla kur, sandbox açık, non-root, audit log aktif."])
cover(prs, "KamuKOD 210", "Kendi Ajanını\nKurmak İster misin?", "4 Hafta · 8 Saat · Çevrimiçi · Türkçe")

OUT_H = f"{OUT}/hermes-intro.pptx"
prs.save(OUT_H)
print(f"  ✓ {OUT_H}")

# === B2B USE CASES DECK ===
print("[2/3] B2B Use Cases Deck...")
prs = Presentation()
prs.slide_width = W; prs.slide_height = H

cover(prs, "Kurumsal Kullanım", "Yapay Zeka Ajanları\nB2B Kullanım Vakaları", "OpenClaw & Hermes ile otonom iş gücü")
title_body(prs, "Finans", "Mevzuat Takip Ajanı",
    ["Sorun: 3 kişilik ekip haftalık SPK/BDDK/TCMB düzenlemelerini tarıyor. 120+ sayfa/hafta.",
     "Çözüm: Hermes ajanı mevzuatı otomatik tarar, ilgili maddeleri işaretler, Slack'ten gönderir.",
     "Sonuç: Haftada 15 insan-saat tasarruf. Gözden kaçan düzenleme riski sıfır."])
title_body(prs, "İK", "CV Ön Eleme Asistanı",
    ["Sorun: Pozisyon başına 300+ başvuru, İK ekibi CV okumaktan işe odaklanamıyor.",
     "Çözüm: Hermes ajanı CV'leri kriterlere göre puanlar, uygun adayları sıralar.",
     "Sonuç: Ön eleme 3 günden 2 saate indi. Yanlış eleme %15'ten %3'e düştü."])
title_body(prs, "BT Operasyon", "Incident Response Ajanı",
    ["Sorun: Gece 3'te gelen sunucu alarmına 30 dk içinde müdahale gerekiyor.",
     "Çözüm: Hermes alarmı alır, log'ları tarar, kök neden analizi yapar, çözer veya bildirir.",
     "Sonuç: Müdahale süresi 5 dk'ya indi. Kritik olmayan alarmların %70'i insansız çözülüyor."])
title_body(prs, "Hukuk", "Sözleşme İnceleme",
    ["Sorun: Her sözleşme 2-3 saat avukat mesaisi. Standart maddeler tekrar tekrar okunuyor.",
     "Çözüm: Ajan sözleşmeyi tarar, riskli maddeleri işaretler, avukata özet sunar.",
     "Sonuç: Sözleşme başına 2 saat tasarruf. Avukatlar sadece riskli maddelere odaklanıyor."])
title_body(prs, "Satış", "Müşteri İçgörü Raporu",
    ["Sorun: Satış görüşmesi öncesi müşteri hakkında bilgi toplamak 45 dk sürüyor.",
     "Çözüm: Ajan CRM, haberler, LinkedIn tarar, tek sayfalık brifing hazırlar, WhatsApp'tan gönderir.",
     "Sonuç: Görüşme hazırlığı 45 dk'dan 2 dk'ya. Kapanan anlaşma oranı %22 arttı."])
title_body(prs, "Kamu", "Vatandaş Başvuru Ön İşleme",
    ["Sorun: Günlük 500+ başvuru manuel sınıflandırılıp yönlendiriliyor.",
     "Çözüm: Ajan başvuruyu okur, sınıflandırır, eksik evrak varsa bilgilendirir, yönlendirir.",
     "Sonuç: İşlem süresi 3 günden 4 saate indi. Vatandaş memnuniyeti arttı."])
cover(prs, "Sonraki Adım", "Kurumunuza Özel\nÇözüm İçin", "github.com/murataslan1/kamuyz-agent-starter")

OUT_B = f"{OUT}/b2b-use-cases.pptx"
prs.save(OUT_B)
print(f"  ✓ {OUT_B}")

# === KAMUKOD 210 DECK ===
print("[3/3] KamuKOD 210 Deck...")
prs = Presentation()
prs.slide_width = W; prs.slide_height = H

cover(prs, "KamuKOD Atölye 210", "OpenClaw & Hermes\nile YZ Ajanları", "4 Hafta · 8 Saat · Çevrimiçi · Türkçe")
title_body(prs, "Problem", "Ekibiniz tekrarlayan\n işlere boğuldu",
    ["Finans: Mevzuat taraması — 15 insan-saat/hafta",
     "İK: CV ön eleme — 3 gün/pozisyon",
     "BT: Gece alarm müdahalesi — 30 dk ortalama",
     "Hukuk: Sözleşme inceleme — 3 saat/belge"])
title_body(prs, "Çözüm", "ChatGPT yetmez.\nOtonom ajan gerek.",
    ["OpenClaw: 385K GitHub yıldızı, 50+ platforma bağlanan iletişim geçidi",
     "Hermes Agent: Kendi kendine öğrenen, zamanla kuruma özel hale gelen ajan",
     "",
     "\"3 saatte 12 Jira ticket'ı kapatan ajan yazdım\" — X kullanıcısı"])
compare_slide(prs, "Neden KamuKOD?", "Kendi başına vs KamuKOD",
    "Kendi başına", ["Kurulum ~4 saat", "Güvenlik kapalı gelir", "%12 zararlı skill riski", "İngilizce doküman", "20-40 saat toplam"],
    "KamuKOD 210", ["İlk oturumda çalışır ✓", "Sandbox açık ✓", "Güvenli kaynaklar ✓", "Türkçe, adım adım ✓", "8 saat ✓"])
title_body(prs, "Program", "4 Oturum, 8 Saat",
    ["1. Kurulum & Telegram bot — Çalışan ajan",
     "2. WhatsApp + Kimlik + Skill — Çok kanallı asistan",
     "3. Google Workspace + Cron — Kurumsal entegrasyon",
     "4. Hermes Learning Loop + Güvenlik — Öğrenen, güvenli sistem"])
cover(prs, "Kayıt", "kamukod.lovable.app\n/atolye/210", "Kontenjan sınırlı")

OUT_K = f"{OUT}/kamukod-210.pptx"
prs.save(OUT_K)
print(f"  ✓ {OUT_K}")

print(f"\n✓ Tüm PPTX'ler {OUT}/ altında")
